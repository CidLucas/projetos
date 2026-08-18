#!/usr/bin/env python3
"""Gera o deck institucional da Deep Blue (v2 — identidade da landing, clara).

Identidade brand-hub (deepblue.company): canvas #F2F2F0, cards brancos borda
#E2E4E0, ink #101828, accent #1D4ED8, kickers JetBrains Mono, termo-chave em
Instrument Serif itálico. Card Consulting dark #101828 (contraste do site).

Uso: python3 build_deck.py [output.pptx]
"""
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------- Design tokens (identidade landing — brand-hub) ----------
CANVAS    = "F2F2F0"
CARD      = "FFFFFF"
CARD_DARK = "101828"
DARK_INK  = "0B1120"
LINE      = "E2E4E0"
DARK_LINE = "1E293B"
INK       = "101828"
MUT       = "5B6472"
MUT_DARK  = "B8BEC9"
AC        = "1D4ED8"
AC2       = "60A5FA"
TAG_BG    = "EDF1FC"   # rgba(29,78,216,.08) sobre branco
PILL_LINE = "D6DEEF"   # azul claro sutil

FONT_T    = "Plus Jakarta Sans"
FONT_B    = "Inter"
FONT_SERIF = "Instrument Serif"
FONT_MONO = "JetBrains Mono"

def rgb(c):
    return RGBColor.from_string(c)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=0, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para_specs in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        for text, style in para_specs:
            r = p.add_run()
            r.text = text
            f = r.font
            f.size = Pt(style.get("size", 14))
            f.bold = style.get("bold", False)
            f.italic = style.get("italic", False)
            f.name = style.get("font", FONT_B)
            f.color.rgb = rgb(style.get("color", INK))
    return tb

def kicker(slide, text, x=0.7, y=0.72):
    """Kicker da landing: JetBrains Mono, uppercase, accent."""
    return add_text(slide, x, y, 11.9, 0.35, [[(text.upper(),
            {"size": 11, "bold": True, "color": AC, "font": FONT_MONO})]])

def title(slide, main, soft=None, x=0.7, y=1.02, size=29, w=11.9,
          line_spacing=1.05):
    """Título PJS 800 ink + termo-chave em Instrument Serif itálico accent."""
    runs = [[(main, {"size": size, "bold": True, "color": INK, "font": FONT_T})]]
    if soft:
        runs[0].append((" " + soft, {"size": size, "bold": False, "italic": True,
                                      "color": AC, "font": FONT_SERIF}))
    return add_text(slide, x, y, w, 1.6, runs, line_spacing=line_spacing)

def sub(slide, text, x=0.7, y=2.45, size=14, w=11.9, color=MUT, line_spacing=1.32):
    return add_text(slide, x, y, w, 0.85, [[(text, {"size": size, "color": color})]],
                    line_spacing=line_spacing)

def card(slide, x, y, w, h, fill=CARD, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    return add_rect(slide, x, y, w, h, fill=fill, line=line, shape=shape)

def tag(slide, text, x, y, fill=TAG_BG, color=AC, size=10, w=None):
    """Pill de tag da landing (mono uppercase)."""
    tw = w if w else 0.32 + 0.088 * len(text)
    sp = add_rect(slide, x, y, tw, 0.34, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x + 0.12, y + 0.075, tw - 0.24, 0.22,
             [[(text.upper(), {"size": size, "bold": True, "color": color,
                               "font": FONT_MONO})]])
    return sp

def footer(slide, page, total=12):
    add_rect(slide, 0.7, 7.08, 11.93, 0.012, fill=LINE)
    add_text(slide, 0.7, 7.15, 6, 0.3, [[("Deep Blue · deepblue.company",
            {"size": 9.5, "color": MUT})]])
    add_text(slide, 11.5, 7.15, 1.13, 0.3, [[(f"{page:02d} / {total:02d}",
            {"size": 9.5, "color": MUT})]], align=PP_ALIGN.RIGHT)

def brand(slide, x=0.7, y=0.5, scale=1.0, dark_text=False):
    """Logo da landing: círculos concêntricos (favicon) + wordmark Serif."""
    colors = ["D0D2D5", "AFB3BB", "5D7FDF", "1D4ED8"]
    steps = [1.0, 0.73, 0.45, 0.18]
    base = 0.34 * scale
    for color, rel in zip(colors, steps):
        d = base * 2 * rel
        add_rect(slide, x + (base - d/2), y + (base - d/2), d, d,
                 fill=color, shape=MSO_SHAPE.OVAL)
    add_text(slide, x + base * 2 + 0.14, y + 0.02, 3.4, 0.4,
             [[("Deep ", {"size": 16.5, "bold": False, "color": INK, "font": FONT_SERIF}),
               ("Blue", {"size": 16.5, "italic": True, "color": MUT, "font": FONT_SERIF})]])

def rings(slide, cx, cy, base=3.0):
    """Círculos concêntricos grandes e quietos (fundo claro)."""
    colors = ["E4E6EA", "D0D6E2", "B3C2E4", "8FABE0"]
    steps = [1.0, 0.73, 0.45, 0.18]
    for color, rel in zip(colors, steps):
        d = base * 2 * rel
        add_rect(slide, cx - d/2, cy - d/2, d, d, fill=color, shape=MSO_SHAPE.OVAL)

def dots(slide, x, y, n=3, color=AC, d=0.07, gap=0.12):
    for i in range(n):
        add_rect(slide, x + i*gap, y, d, d, fill=color, shape=MSO_SHAPE.OVAL)

def blk_mut(slide, text, x, y, w=5.2, size=13.5, color=MUT, line_spacing=1.3):
    return add_text(slide, x, y, w, 1.6, [[(text, {"size": size, "color": color})]],
                    line_spacing=line_spacing)

# ---------- Build ----------
def build(path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    # ===== Slide 1 — Capa =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    rings(s, 11.2, 3.9, base=3.1)
    brand(s, 0.7, 0.62, scale=1.15)
    add_rect(s, 0.7, 2.5, 0.55, 0.045, fill=AC)
    add_text(s, 0.7, 2.75, 7.6, 0.4, [[("INTELIGÊNCIA ARTIFICIAL COMO FERRAMENTA DE NEGÓCIO",
            {"size": 12, "bold": True, "color": AC, "font": FONT_MONO})]])
    add_text(s, 0.7, 3.25, 8.6, 2.2,
             [[("Sua operação já funciona.", {"size": 40, "bold": True, "color": INK, "font": FONT_T})],
              [("A IA pode ser o que falta para ela ", {"size": 40, "bold": True, "color": INK, "font": FONT_T}),
               ("render mais.", {"size": 40, "bold": False, "italic": True, "color": AC, "font": FONT_SERIF})]],
             line_spacing=1.02, space_after=6)
    add_text(s, 0.7, 5.7, 7.8, 0.9,
             [[("Entregamos sistemas que agregam valor, com resultado mensurável. ", {"size": 15, "color": MUT}),
               ("Consultoria primeiro. Tecnologia depois.", {"size": 15, "bold": True, "color": INK})]],
             line_spacing=1.35)
    add_text(s, 0.7, 6.85, 8, 0.35, [[("deepblue.company", {"size": 12, "color": MUT})]])
    s.notes_slide.notes_text_frame.text = (
        "Abertura. Uma frase: a Deep Blue existe para PMEs que já rodam bem e querem render mais. "
        "Nosso ponto de partida é sempre o negócio; a IA entra depois, como ferramenta.")

    # ===== Slide 2 — Quem somos =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "A empresa")
    title(s, "A Deep Blue ajuda empresas a usar IA como", soft="ferramenta de negócio.", size=27)
    sub(s, "Fundada para um perfil específico: PMEs que já rodam bem e querem render mais. "
           "Duas frentes, uma lógica — entender o negócio antes de propor tecnologia.", y=2.45)
    # Card Labs (branco)
    card(s, 0.7, 3.45, 5.8, 2.75)
    tag(s, "Labs · assinatura", 1.0, 3.75)
    add_text(s, 1.0, 4.25, 5.2, 0.5, [[("Deep Blue Labs", {"size": 19, "bold": True, "color": INK, "font": FONT_T})]])
    blk_mut(s, "Produtos em assinatura. Plataforma Blu, agentes de IA e assistentes que "
               "centralizam conhecimento e automatizam o dia a dia.", 1.0, 4.85, size=13.5)
    # Card Consulting (dark — contraste do site)
    card(s, 6.83, 3.45, 5.8, 2.75, fill=CARD_DARK, line=DARK_LINE)
    tag(s, "Consulting · escopo fechado", 7.13, 3.75, color=AC2, fill="1B2543")
    add_text(s, 7.13, 4.25, 5.2, 0.5, [[("Deep Blue Consulting", {"size": 19, "bold": True, "color": "F2F2F0", "font": FONT_T})]])
    blk_mut(s, "Projetos de diagnóstico e estratégia. Mapeamos processos, identificamos onde "
               "a IA agrega valor e priorizamos casos de uso.", 7.13, 4.85, color=MUT_DARK, size=13.5)
    footer(s, 2)
    s.notes_slide.notes_text_frame.text = (
        "Apresentar as duas frentes como complementares: Labs é o produto contínuo, "
        "Consulting é a porta de entrada. Muitos clientes começam no diagnóstico e viram assinante depois.")

    # ===== Slide 3 — O contexto =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "O contexto")
    title(s, "PMEs rodam o país. E o dono decide", soft="sozinho.", size=29)
    sub(s, "Dados de mercado que explicam por que a IA vira ferramenta — e não brinquedo.", y=2.42)
    stats = [
        ("90%", "dos líderes de PMEs tomam decisões sozinhos em 5 áreas diferentes, sem apoio estruturado", "Itaú · Locomotiva"),
        ("30%", "do PIB brasileiro é movido por PMEs, que respondem por metade dos empregos ativos", "Sebrae"),
        ("9 em 10", "PMEs relatam dificuldade na gestão financeira do dia a dia", "Pesquisa Sebrae 2024"),
    ]
    for i, (num, txt, src) in enumerate(stats):
        x = 0.7 + i * 4.05
        card(s, x, 3.3, 3.75, 3.05)
        add_text(s, x + 0.3, 3.62, 3.15, 1.0, [[(num, {"size": 44, "bold": True, "color": AC, "font": FONT_T})]])
        add_text(s, x + 0.3, 4.6, 3.15, 1.35, [[(txt, {"size": 12.5, "color": MUT})]], line_spacing=1.3)
        add_text(s, x + 0.3, 5.9, 3.15, 0.3, [[(src, {"size": 10, "color": "8A93A3", "italic": True})]])
    footer(s, 3)
    s.notes_slide.notes_text_frame.text = (
        "Fonte: Itaú/Locomotiva, Sebrae (2024-2025). O dono de PME decide sozinho em 5 áreas; "
        "a Deep Blue existe pra tirar peso das costas dele.")

    # ===== Slide 4 — Método =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Método")
    title(s, "Não começamos pela tecnologia. Começamos pelo", soft="negócio.", size=28)
    steps = [
        ("01", "Escutar", "Entrevistas com as pessoas, não só com os dados. As dores aparecem no dia a dia, não no relatório."),
        ("02", "Entender", "Processos, decisões, gargalos. Onde o tempo some e onde o conhecimento fica preso."),
        ("03", "Mapear", "Oportunidades priorizadas por impacto e viabilidade. O diagnóstico vira um mapa, não um laudo."),
    ]
    for i, (num, t, txt) in enumerate(steps):
        x = 0.7 + i * 4.05
        card(s, x, 3.3, 3.75, 3.05)
        add_text(s, x + 0.3, 3.62, 3.15, 0.6, [[(num, {"size": 14, "bold": True, "color": "9AA3B2", "font": FONT_MONO})]])
        add_text(s, x + 0.3, 4.2, 3.15, 0.5, [[(t, {"size": 20, "bold": True, "color": INK, "font": FONT_T})]])
        blk_mut(s, txt, x + 0.3, 4.85, w=3.15, size=12.5)
    add_text(s, 0.7, 6.55, 11.9, 0.4,
             [[("Todo projeto de IA começa com pessoas: escutamos as dores do cliente, entendemos o processo e mapeamos oportunidades.",
                {"size": 13, "italic": True, "color": INK})]])
    footer(s, 4)
    s.notes_slide.notes_text_frame.text = (
        "Frase do fundador (padrão ouro). Consultoria primeiro, tecnologia depois. "
        "Isso diferencia a Deep Blue de quem vende ferramenta pronta sem entender o negócio.")

    # ===== Slide 5 — O diagnóstico =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "O que encontramos")
    title(s, "As dores se repetem. Conhecimento preso, decisão", soft="solitária,", size=26)
    add_text(s, 0.7, 1.78, 11.9, 0.5, [[("processo manual.", {"size": 26, "bold": True, "color": INK, "font": FONT_T})]])
    pains = [
        ("Conhecimento espalhado", "Documentos, planilhas e SOPs em lugares diferentes. O que a equipe sabe não vira ativo da empresa."),
        ("Decisão solitária", "O dono decide sozinho, sem ter os números na mão no momento da decisão. Tudo depende da cabeça dele."),
        ("Trabalho manual", "O tempo que poderia gerar negócio é gasto repetindo o mesmo processo. Cada proposta, cada relatório, do zero."),
    ]
    for i, (t, txt) in enumerate(pains):
        x = 0.7 + i * 4.05
        card(s, x, 2.75, 3.75, 3.6)
        dots(s, x + 0.3, 3.05, n=3, d=0.075, gap=0.13)
        add_text(s, x + 0.3, 3.4, 3.15, 0.5, [[(t, {"size": 17, "bold": True, "color": INK, "font": FONT_T})]])
        blk_mut(s, txt, x + 0.3, 4.05, w=3.15, size=12.5)
    footer(s, 5)
    s.notes_slide.notes_text_frame.text = (
        "Essas três dores aparecem em quase todo diagnóstico. Use exemplos do cliente real "
        "que você está apresentando: troque os genéricos pelas dores específicas dele.")

    # ===== Slide 6 — Capabilities =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Capabilities")
    title(s, "Ferramentas de IA a serviço do seu", soft="negócio.", size=29)
    sub(s, "Três frentes. Todas inseridas nos pontos de impacto que o diagnóstico revelou.", y=2.42)
    caps = [
        ("Plataforma", "Sistemas completos: documentos, conhecimento e agentes de IA num ambiente só. A operação centralizada."),
        ("Fluxos de agentes", "Processos automatizados com agentes integrados aos sistemas que o cliente já usa. Cada um na sua função."),
        ("Assistente diário", "Um agente que entende o negócio e responde em linguagem natural. Perguntou, resolveu — com contexto."),
    ]
    for i, (t, txt) in enumerate(caps):
        x = 0.7 + i * 4.05
        card(s, x, 3.3, 3.75, 3.05)
        tag(s, t, x + 0.3, 3.62, w=1.85)
        blk_mut(s, txt, x + 0.3, 4.3, w=3.15, size=12.5)
        add_rect(s, x + 0.3, 5.95, 0.5, 0.045, fill=AC)
    footer(s, 6)
    s.notes_slide.notes_text_frame.text = (
        "Visão geral antes de aprofundar. As três frentes respondem a três perguntas do cliente: "
        "'onde fica minha operação?', 'quem executa meu processo?', 'quem responde minhas perguntas?'")

    # ===== Slide 7 — Plataforma Blu =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Capability · Plataforma")
    title(s, "Plataforma Blu. O escritório virtual da sua", soft="empresa.", size=28)
    sub(s, "Um ambiente só para o que hoje está espalhado: documentos, conhecimento e agentes.", y=2.42)
    feats = [
        ("Agentes de IA", "Trabalham nos seus processos: atendimento, documentos, relatórios, rotinas."),
        ("Pipeline de documentos", "Do lead ao contrato, tudo num fluxo único e rastreável."),
        ("Sala de estratégia", "Busca semântica no conhecimento da empresa. Pergunta em linguagem natural, resposta com contexto."),
        ("Gestão de conhecimento", "O que a equipe sabe vira ativo da empresa — acessível, não preso na cabeça de ninguém."),
    ]
    for i, (t, txt) in enumerate(feats):
        col, row = i % 2, i // 2
        x = 0.7 + col * 6.03
        y = 3.3 + row * 1.8
        card(s, x, y, 5.8, 1.58)
        add_text(s, x + 0.3, y + 0.2, 5.2, 0.4, [[(t, {"size": 16, "bold": True, "color": INK, "font": FONT_T})]])
        blk_mut(s, txt, x + 0.3, y + 0.66, size=12.5)
    footer(s, 7)
    s.notes_slide.notes_text_frame.text = (
        "Blu é o carro-chefe do Labs. Se o cliente tem a operação espalhada em planilha, e-mail e Drive, "
        "o Blu centraliza. Demo ao vivo se possível: mostrar busca semântica na sala de estratégia.")

    # ===== Slide 8 — Agentes =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Capability · Agentes")
    title(s, "Agentes que trabalham dentro do seu", soft="processo.", size=28)
    sub(s, "Implantados nos sistemas que o cliente já usa. Cada um com uma função clara.", y=2.42)
    agents = [
        ("TalentFlow", "RH — triagem de currículos, agendamento, onboarding, análise de turnover."),
        ("KnowledgeBase", "Conhecimento — respostas com o contexto da empresa, não texto genérico."),
        ("ProjectPilot", "Projetos — prazos, riscos, status. O PM que não dorme."),
        ("OpsFlow", "Operações — rotinas e fluxos internos automatizados."),
    ]
    for i, (t, txt) in enumerate(agents):
        col, row = i % 2, i // 2
        x = 0.7 + col * 6.03
        y = 3.3 + row * 1.8
        card(s, x, y, 5.8, 1.58)
        add_text(s, x + 0.3, y + 0.2, 5.2, 0.4, [[(t, {"size": 16, "bold": True, "color": INK, "font": FONT_T})]])
        blk_mut(s, txt, x + 0.3, y + 0.66, size=12.5)
    add_text(s, 0.7, 6.55, 11.9, 0.4,
             [[("Nossa operação roda com 6 agentes assim. A mesma tecnologia que entregamos para clientes.",
                {"size": 13, "italic": True, "color": INK})]])
    footer(s, 8)
    s.notes_slide.notes_text_frame.text = (
        "Agentes sob medida são o ticket mais alto (implantação + mensalidade). "
        "Contar que a Deep Blue usa os próprios agentes para rodar a operação: é vitrine, não teoria.")

    # ===== Slide 9 — Assistente diário =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Capability · Assistente")
    title(s, "Pergunte em linguagem natural. Ele entende o seu", soft="negócio.", size=27)
    sub(s, "Um assistente diário — um Hermes trabalhando para você, com o contexto da empresa.", y=2.42)
    card(s, 0.7, 3.3, 7.4, 3.1)
    flow = [
        ("1 · Entende o contexto", "Processo, gargalos e histórico da empresa. Não responde no vácuo."),
        ("2 · Consulta as ferramentas por você", "Busca nos documentos, puxa os números, executa rotinas."),
        ("3 · Responde com contexto", "Resposta com a cara do seu negócio — não texto genérico de IA."),
    ]
    for i, (t, txt) in enumerate(flow):
        y = 3.55 + i * 0.95
        add_text(s, 1.05, y, 6.8, 0.8,
                 [[(t + "  ", {"size": 14, "bold": True, "color": INK, "font": FONT_T}),
                   (txt, {"size": 13, "color": MUT})]], line_spacing=1.25)
    card(s, 8.4, 3.3, 4.23, 3.1, fill="F8F9FC", line=PILL_LINE)
    tag(s, "Exemplo", 8.75, 3.62, fill=TAG_BG, w=1.05)
    add_text(s, 8.75, 4.15, 3.6, 2.1,
             [[("\u201cQuantas propostas foram enviadas este mês e quais estão paradas há mais de uma semana?\u201d",
                {"size": 13, "italic": True, "color": INK, "font": FONT_SERIF})],
              [("→ Busca, cruza e responde com os números da casa.", {"size": 12.5, "color": MUT})]],
             line_spacing=1.35, space_after=10)
    footer(s, 9)
    s.notes_slide.notes_text_frame.text = (
        "O assistente diário é o 'Hermes do cliente': consulta em linguagem natural, entende o negócio, "
        "usa as ferramentas. Demo ideal: fazer a pergunta ao vivo e mostrar a resposta com contexto.")

    # ===== Slide 10 — O método =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Como entra a tecnologia")
    title(s, "Diagnóstico primeiro. Ferramentas nos pontos de", soft="impacto.", size=28)
    sub(s, "Entendemos as dores do cliente e inserimos as ferramentas exatamente onde elas mudam o dia a dia.", y=2.42)
    card(s, 0.7, 3.45, 5.8, 2.7)
    tag(s, "1 · Diagnóstico", 1.0, 3.75, w=2.2)
    blk_mut(s, "Gargalos e dores mapeados com o cliente: onde o tempo some, onde a decisão trava, "
               "onde o conhecimento fica preso.", 1.0, 4.4, size=13.5)
    add_text(s, 6.9, 4.25, 0.9, 0.9, [[("→", {"size": 40, "bold": True, "color": AC, "font": FONT_T})]])
    card(s, 7.83, 3.45, 4.8, 2.7)
    tag(s, "2 · Inserção", 8.13, 3.75, w=2.2)
    blk_mut(s, "Ferramentas de IA nos pontos de impacto do diagnóstico. Cada uma responde a uma dor mapeada.",
            8.13, 4.4, size=13.5)
    add_text(s, 0.7, 6.55, 11.9, 0.4,
             [[("A IA é uma poderosa ferramenta que usamos para impactar o diagnóstico que fizemos.",
                {"size": 13, "italic": True, "color": INK})]])
    footer(s, 10)
    s.notes_slide.notes_text_frame.text = (
        "O coração da proposta: nada de ferramenta sem diagnóstico. Esse é o slide que separa "
        "a Deep Blue de quem vende IA pronta. Feche com a frase do fundador.")

    # ===== Slide 11 — Prova viva =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Prova viva")
    title(s, "Usamos o que", soft="vendemos.", size=30)
    sub(s, "Nossa operação roda com agentes de IA. A mesma tecnologia que entregamos para clientes.", y=2.42)
    team = [
        ("Writer", "Propostas, documentação, conteúdo"),
        ("PM", "Cronogramas, riscos, status"),
        ("Sales", "Leads, follow-ups, pipeline"),
        ("Dev", "Código, code review, testes"),
        ("Legal", "Contratos e compliance"),
        ("Finance", "Fluxo de caixa e margem por projeto"),
    ]
    for i, (t, txt) in enumerate(team):
        col, row = i % 3, i // 3
        x = 0.7 + col * 4.05
        y = 3.4 + row * 1.65
        card(s, x, y, 3.75, 1.45)
        add_text(s, x + 0.3, y + 0.18, 3.15, 0.4, [[("Hermes " + t, {"size": 15, "bold": True, "color": AC, "font": FONT_T})]])
        add_text(s, x + 0.3, y + 0.66, 3.15, 0.7, [[(txt, {"size": 12, "color": MUT})]], line_spacing=1.2)
    footer(s, 11)
    s.notes_slide.notes_text_frame.text = (
        "Os 6 agentes Hermes internos. Quando o cliente perguntar 'isso funciona?', "
        "a resposta é: a gente opera assim. Cada proposta que vocês estão vendo foi escrita por um agente.")

    # ===== Slide 12 — Próximos passos =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, fill=CANVAS)
    rings(s, 11.4, 1.7, base=2.2)
    brand(s, 0.7, 0.5, scale=1.0)
    kicker(s, "Próximos passos")
    title(s, "Conta o seu desafio. A gente mapeia onde a IA", soft="agrega.", size=28)
    steps = [
        ("1", "Diagnóstico gratuito", "Entendemos seu processo e suas dores. Sem compromisso, sem fórmula pronta."),
        ("2", "Mapa de impacto", "Onde a IA agrega valor no seu negócio, com priorização por impacto e viabilidade."),
        ("3", "Proposta", "Ferramentas nos pontos de impacto, com resultado mensurável e cronograma."),
    ]
    for i, (num, t, txt) in enumerate(steps):
        x = 0.7 + i * 4.05
        card(s, x, 3.4, 3.75, 2.7)
        add_text(s, x + 0.3, 3.7, 1, 0.6, [[(num, {"size": 28, "bold": True, "color": "9AA3B2", "font": FONT_MONO})]])
        add_text(s, x + 0.3, 4.35, 3.15, 0.5, [[(t, {"size": 17, "bold": True, "color": INK, "font": FONT_T})]])
        blk_mut(s, txt, x + 0.3, 4.95, w=3.15, size=12.5)
    add_text(s, 0.7, 6.4, 11.9, 0.5,
             [[("deepblue.company  ·  ", {"size": 13, "color": MUT}),
               ("formly.ink", {"size": 13, "color": AC, "bold": True}),
               ("  ·  ", {"size": 13, "color": MUT}),
               ("app.mcp-brain.com", {"size": 13, "color": AC, "bold": True})]])
    footer(s, 12)
    s.notes_slide.notes_text_frame.text = (
        "Fechamento com ação: agendar o diagnóstico gratuito. O CTA real é a reunião de "
        "descoberta. Levar contatos: site, Formly, MCP Brain.")

    prs.save(path)
    print(f"OK -> {path} ({len(prs.slides._sldIdLst)} slides)")

if __name__ == "__main__":
    build(sys.argv[1] if len(sys.argv) > 1 else "deck-empresa-v2.pptx")
